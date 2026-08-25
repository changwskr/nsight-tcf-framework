# -*- coding: utf-8 -*-
"""12.png → editable PPTX (text / shapes / tables / images)."""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
SRC = ROOT / "원본" / "12.png"
OUT = ROOT / "마케팅플랫폼_표준_아키텍처_구성.pptx"
ASSETS = ROOT / "_pptx_assets"
FONT = "맑은 고딕"

NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LT = RGBColor(0xBD, 0xD7, 0xEE)
GREEN = RGBColor(0x70, 0xAD, 0x47)
GREEN_DK = RGBColor(0x54, 0x8B, 0x54)
LIME = RGBColor(0xA9, 0xD0, 0x8E)
GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_BOX = RGBColor(0xD9, 0xD9, 0xD9)
GRAY_DK = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xC0, 0x00, 0x00)
IF_FILL = RGBColor(0xDE, 0xEA, 0xF6)
LINE = RGBColor(0x7F, 0x7F, 0x7F)


def set_run_font(run, size, color=BLACK, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", FONT)


def set_shape_text(shape, text, size=10, color=BLACK, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    shape.text_frame.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size, color, bold)
    try:
        tf._bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass


def fill_line(shape, fill, line=LINE, line_pt=0.75):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_pt)


def add_box(slide, l, t, w, h, text, fill=BLUE, font=9, color=WHITE, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE):
    sh = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_line(sh, fill, line)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = 0.12
        except Exception:
            pass
    set_shape_text(sh, text, font, color, bold)
    return sh


def add_label(slide, l, t, w, h, text, size=9, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    set_shape_text(tb, text, size, color, bold, align=align, anchor=MSO_ANCHOR.TOP)
    return tb


def add_layer_tag(slide, l, t, w, h, text):
    return add_box(slide, l, t, w, h, text, fill=NAVY, font=8, color=WHITE, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def add_arrow_h(slide, l, t, w=0.22, h=0.12):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_line(sh, GRAY_DK, GRAY_DK, 0.5)
    return sh


def add_arrow_v(slide, l, t, w=0.12, h=0.18):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_line(sh, GRAY_DK, GRAY_DK, 0.5)
    return sh


def add_hline(slide, l, t, w, color=LINE, pt=1.0, dash=False):
    if not dash:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(pt))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        return sh
    x = l
    while x < l + w:
        seg_w = min(0.08, l + w - x)
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(t), Inches(seg_w), Pt(pt))
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.fill.background()
        x += 0.14
    return None


def style_table(table, header=True):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run_font(run, 9 if r or not header else 10, WHITE if header and r == 0 else BLACK, bold=(header and r == 0) or c == 0)
            fill = cell.fill
            fill.solid()
            if header and r == 0:
                fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
            else:
                fill.fore_color.rgb = WHITE


def prepare_assets():
    ASSETS.mkdir(exist_ok=True)
    im = Image.open(SRC).convert("RGB")
    # content crop (screen photo has dark frame)
    crop = im.crop((400, 70, 1620, 960))
    content_path = ASSETS / "12_content.png"
    crop.save(content_path)
    # decorative NH / SK badges as simple images
    for name, bg, fg, text in [
        ("nh.png", (0, 102, 51), (255, 255, 255), "농협중앙회"),
        ("sk.png", (220, 30, 40), (255, 255, 255), "SK 주식회사"),
    ]:
        badge = Image.new("RGB", (420, 70), bg)
        badge.save(ASSETS / name)  # text drawn in PPT instead; keep color bar
    return content_path


def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title + accent bar
    add_label(slide, 0.35, 0.12, 10, 0.38, "마케팅플랫폼 표준 아키텍처 구성", 22, BLACK, True)
    add_box(slide, 0.35, 0.48, 0.35, 0.06, "", fill=BLUE, shape=MSO_SHAPE.RECTANGLE, line=BLUE)
    add_box(slide, 0.70, 0.48, 0.35, 0.06, "", fill=LIME, shape=MSO_SHAPE.RECTANGLE, line=LIME)
    add_box(slide, 1.05, 0.48, 11.9, 0.06, "", fill=GREEN_DK, shape=MSO_SHAPE.RECTANGLE, line=GREEN_DK)

    # intro
    intro = (
        "마케팅플랫폼은 IaaS 기반 아키텍처로, 실시간 고객 반응형 서비스 및 Customer Single View를 "
        "주 제공하는 시스템으로 단말관리 및 배포를 위한 UI 영역과 NH Cloud Framework 영역으로 구현합니다."
    )
    intro_box = add_box(
        slide, 0.35, 0.60, 12.6, 0.42, intro,
        fill=RGBColor(0xFF, 0xFB, 0xF0), font=10, color=GRAY_DK, bold=False,
        shape=MSO_SHAPE.RECTANGLE, line=GREEN_DK,
    )
    set_shape_text(intro_box, intro, 10, GRAY_DK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    # ---- Client Access ----
    add_layer_tag(slide, 0.20, 1.15, 0.95, 0.55, "Client\nAccess")
    add_box(slide, 1.30, 1.12, 1.35, 0.32, "전용 Browser\n(마케팅플랫폼)", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_box(slide, 2.85, 1.12, 1.10, 0.32, "통합단말", fill=GRAY_BOX, font=9, color=BLACK, bold=False)
    add_box(slide, 1.30, 1.55, 0.70, 0.28, "L4", fill=GRAY_BOX, font=9, color=BLACK, bold=True)
    add_box(slide, 2.15, 1.55, 0.70, 0.28, "L4", fill=GRAY_BOX, font=9, color=BLACK, bold=True)
    add_box(slide, 3.00, 1.55, 0.80, 0.28, "MCA", fill=GRAY_BOX, font=9, color=BLACK, bold=True)
    add_label(slide, 2.85, 1.42, 0.55, 0.16, "싱글뷰", 8, RED, True, PP_ALIGN.CENTER)
    add_arrow_v(slide, 1.55, 1.44, 0.12, 0.12)
    add_arrow_v(slide, 3.30, 1.44, 0.12, 0.12)
    add_arrow_h(slide, 2.85, 1.62, 0.15, 0.12)

    # ---- Service ----
    add_layer_tag(slide, 0.20, 2.55, 0.95, 0.70, "Service")
    # left IaaS UI
    left = add_box(slide, 1.25, 1.95, 3.55, 2.05, "", fill=GRAY, font=8, color=BLACK, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    set_shape_text(left, "", 8)
    add_label(slide, 1.35, 1.98, 3.3, 0.22, "IaaS ( UI 서비스플랫폼 )", 10, NAVY, True)
    add_box(slide, 2.35, 2.22, 1.35, 0.28, "WEB", fill=BLUE, font=10)
    add_arrow_v(slide, 2.95, 2.50, 0.12, 0.14)
    node_l = add_box(slide, 1.45, 2.68, 3.15, 1.15, "", fill=WHITE, font=8, color=BLACK, bold=False)
    set_shape_text(node_l, "", 8)
    add_label(slide, 1.55, 2.70, 1.5, 0.18, "Service Node", 8, GRAY_DK, True)
    add_box(slide, 1.60, 2.95, 1.30, 0.55, "단말관리\nWAS", fill=BLUE, font=10)
    add_box(slide, 3.05, 2.95, 1.30, 0.55, "단말배포\nWAS", fill=BLUE, font=10)
    add_label(slide, 1.55, 3.55, 2.8, 0.20, "WebTopSuite", 8, GRAY_DK, False, PP_ALIGN.CENTER)

    # center bridge
    add_box(slide, 4.95, 2.35, 0.85, 0.32, "배치 AP", fill=BLUE, font=9)
    add_box(slide, 4.95, 2.80, 0.85, 0.32, "RD", fill=BLUE, font=9)
    add_box(slide, 4.95, 3.25, 0.85, 0.32, "ETL", fill=BLUE, font=9)

    # right IaaS marketing
    right = add_box(slide, 5.95, 1.95, 6.95, 2.05, "", fill=GRAY, font=8, color=BLACK, bold=False)
    set_shape_text(right, "", 8)
    add_label(slide, 6.05, 1.98, 4.5, 0.22, "IaaS ( 마케팅플랫폼 )  ·  NH Cloud Framework", 10, NAVY, True)
    add_box(slide, 7.55, 2.22, 1.50, 0.28, "WEB", fill=BLUE, font=10)
    add_arrow_v(slide, 8.20, 2.50, 0.12, 0.14)

    svc = add_box(slide, 6.10, 2.68, 3.55, 0.70, "", fill=WHITE, font=8, color=BLACK, bold=False)
    set_shape_text(svc, "", 8)
    add_label(slide, 6.15, 2.70, 1.4, 0.16, "Service Node", 8, GRAY_DK, True)
    add_box(slide, 6.20, 2.90, 1.60, 0.40, "[마케팅플랫폼]\nWAS", fill=BLUE, font=8)
    add_box(slide, 7.90, 2.90, 1.60, 0.40, "[마케팅허브]\nWAS", fill=BLUE, font=8)

    # IMDG
    imdg_wrap = add_box(slide, 9.85, 2.68, 1.45, 0.70, "", fill=WHITE, font=8, color=BLACK, bold=False, line=GREEN)
    set_shape_text(imdg_wrap, "", 8)
    add_label(slide, 9.90, 2.70, 1.35, 0.16, "공통인프라 Node", 7, GRAY_DK, True, PP_ALIGN.CENTER)
    add_box(slide, 10.05, 2.90, 1.05, 0.38, "IMDG", fill=GREEN, font=11)
    add_label(slide, 11.35, 2.75, 1.4, 0.35, "Session/\nData Cache", 8, RED, True)

    # realtime
    rt = add_box(slide, 6.10, 3.45, 6.55, 0.48, "", fill=WHITE, font=8, color=BLACK, bold=False)
    set_shape_text(rt, "", 8)
    add_label(slide, 6.15, 3.46, 1.8, 0.16, "실시간 이벤트 처리", 8, GRAY_DK, True)
    add_box(slide, 6.20, 3.62, 2.05, 0.26, "[EBM] 고객행동데이터 Kafka", fill=BLUE, font=7)
    add_box(slide, 8.35, 3.62, 1.85, 0.26, "[EBM] 행동처리", fill=BLUE, font=8)
    add_box(slide, 10.30, 3.62, 2.15, 0.26, "[EBM] 실시간처리 EBM", fill=BLUE, font=8)

    add_hline(slide, 1.25, 4.10, 11.7, LINE, 1.0, dash=True)

    # ---- I/F ----
    add_layer_tag(slide, 0.20, 4.25, 0.95, 0.45, "I/F")
    if_items = [
        "SSO", "Cruz APIM", "EAI", "FOS", "MFT", "ETL", "CDC",
        "대내 MCA", "대외 MCA", "GSE", "통합관제", "IT메타", "Nexus",
        "Wily Collector", "배치작업관리",
    ]
    x0, y0, bw, bh, gap = 1.25, 4.20, 0.52, 0.55, 0.03
    for i, name in enumerate(if_items):
        add_box(slide, x0 + i * (bw + gap), y0, bw, bh, name, fill=IF_FILL, font=6, color=BLACK, bold=False)

    # right external grid
    grid = [
        ["연계포털", "빅데이터", "계정계", "대외기관"],
        ["마이데이터", "올원뱅크", "디지털플랫폼", "타금융"],
    ]
    gx, gy = 9.70, 4.18
    for r, row in enumerate(grid):
        for c, name in enumerate(row):
            add_box(slide, gx + c * 0.85, gy + r * 0.28, 0.81, 0.26, name, fill=IF_FILL, font=7, color=BLACK, bold=False)

    add_hline(slide, 1.25, 4.78, 11.7, LINE, 1.0, dash=True)

    # ---- Data ----
    add_layer_tag(slide, 0.20, 5.00, 0.95, 0.40, "Data")
    add_box(slide, 1.40, 4.90, 1.35, 0.55, "단말관리 DB", fill=GRAY_DK, font=9, color=WHITE, shape=MSO_SHAPE.CAN)
    exa = add_box(slide, 3.00, 4.88, 2.60, 0.60, "", fill=GRAY, font=8, color=BLACK, bold=False)
    set_shape_text(exa, "", 8)
    add_box(slide, 3.15, 4.93, 1.05, 0.38, "ADW", fill=GRAY_DK, font=10, color=WHITE, shape=MSO_SHAPE.CAN)
    add_box(slide, 4.35, 4.93, 1.05, 0.38, "RDW", fill=GRAY_DK, font=10, color=WHITE, shape=MSO_SHAPE.CAN)
    add_label(slide, 3.10, 5.30, 2.4, 0.18, "Oracle Exa", 8, GRAY_DK, True, PP_ALIGN.CENTER)
    add_box(slide, 5.90, 4.90, 1.55, 0.55, "마케팅플랫폼", fill=GRAY_DK, font=9, color=WHITE, shape=MSO_SHAPE.CAN)

    add_hline(slide, 1.25, 5.55, 11.7, LINE, 1.0, dash=True)

    # ---- 형상배포 ----
    add_layer_tag(slide, 0.20, 6.05, 0.95, 0.55, "형상배포")
    # Dev path
    add_label(slide, 1.30, 5.62, 0.9, 0.16, "GitLab", 7, GRAY_DK, False, PP_ALIGN.CENTER)
    add_box(slide, 1.25, 5.78, 1.05, 0.32, "코드버전관리", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_label(slide, 2.35, 5.95, 0.7, 0.16, "dev merge", 7, RED, True, PP_ALIGN.CENTER)
    add_arrow_h(slide, 2.35, 5.85)
    add_label(slide, 2.95, 5.62, 1.0, 0.16, "GitLabRunner", 7, GRAY_DK, False, PP_ALIGN.CENTER)
    add_box(slide, 2.90, 5.78, 0.85, 0.32, "빌드", fill=GRAY_BOX, font=9, color=BLACK, bold=False)
    add_arrow_h(slide, 3.80, 5.85)
    add_box(slide, 4.05, 5.78, 1.05, 0.32, "Dev Pipeline", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_box(slide, 4.20, 5.62, 0.85, 0.16, "취약점점검", fill=RGBColor(0xFC, 0xE4, 0xD6), font=7, color=BLACK, bold=False)
    add_arrow_h(slide, 5.15, 5.85)
    add_box(slide, 5.40, 5.78, 0.95, 0.32, "개발환경", fill=GRAY_BOX, font=9, color=BLACK, bold=False)

    # Master merge note
    add_label(slide, 6.40, 5.62, 2.2, 0.16, "Master merge (통제점검)", 8, RED, True)

    # Prod upper
    add_box(slide, 6.50, 5.78, 0.85, 0.28, "eCAMS", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_arrow_h(slide, 7.40, 5.84)
    add_box(slide, 7.65, 5.78, 0.70, 0.28, "빌드", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_box(slide, 7.55, 5.62, 0.85, 0.16, "취약점점검", fill=RGBColor(0xFC, 0xE4, 0xD6), font=7, color=BLACK, bold=False)
    add_arrow_h(slide, 8.40, 5.84)
    add_box(slide, 8.65, 5.78, 0.95, 0.28, "운영환경", fill=GRAY_BOX, font=8, color=BLACK, bold=False)

    # Prod lower
    add_label(slide, 6.50, 6.22, 1.0, 0.14, "GitLabRunner", 7, GRAY_DK, False, PP_ALIGN.CENTER)
    add_box(slide, 6.50, 6.35, 0.70, 0.28, "빌드", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_arrow_h(slide, 7.25, 6.42)
    add_box(slide, 7.50, 6.35, 1.15, 0.28, "운영-Pipeline", fill=GRAY_BOX, font=8, color=BLACK, bold=False)
    add_arrow_h(slide, 8.70, 6.42)
    add_box(slide, 8.95, 6.35, 0.95, 0.28, "운영환경", fill=GRAY_BOX, font=8, color=BLACK, bold=False)

    # footer
    add_hline(slide, 0.35, 6.80, 12.6, LINE, 0.75)
    add_label(slide, 0.40, 6.90, 2.5, 0.28, "농협중앙회", 11, GREEN_DK, True)
    add_label(slide, 6.0, 6.90, 1.3, 0.28, "- 5 -", 11, BLACK, False, PP_ALIGN.CENTER)
    add_label(slide, 10.8, 6.90, 2.2, 0.28, "SK 주식회사", 11, RED, True, PP_ALIGN.RIGHT)


def build_slide2(prs, content_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_label(slide, 0.35, 0.15, 12, 0.35, "원본 이미지 (12.png · 콘텐츠 영역 크롭)", 18, NAVY, True)
    slide.shapes.add_picture(str(content_path), Inches(0.55), Inches(0.60), width=Inches(12.2))
    add_label(slide, 0.35, 7.05, 12.5, 0.25, "편집용 재구성은 1페이지(도형/텍스트), 구성요소 표는 3페이지를 사용합니다.", 10, GRAY_DK)


def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_label(slide, 0.35, 0.15, 12, 0.35, "계층별 구성요소 표", 20, NAVY, True)

    rows = [
        ("계층", "구성요소", "비고"),
        ("Client Access", "전용 Browser(마케팅플랫폼), 통합단말, L4, MCA", "싱글뷰 연계"),
        ("Service / UI", "WEB, 단말관리 WAS, 단말배포 WAS, WebTopSuite", "IaaS UI 서비스플랫폼"),
        ("Service / 마케팅", "[마케팅플랫폼] WAS, [마케팅허브] WAS, IMDG", "NH Cloud Framework"),
        ("실시간 이벤트", "[EBM] Kafka, 행동처리, 실시간처리 EBM", "고객행동 실시간"),
        ("Bridge", "배치 AP, RD, ETL", "플랫폼 간 연계"),
        ("I/F", "SSO, Cruz APIM, EAI, FOS, MFT, ETL, CDC, MCA, GSE …", "공통 인터페이스"),
        ("Data", "단말관리 DB, ADW, RDW(Oracle Exa), 마케팅플랫폼", "데이터 계층"),
        ("형상배포", "GitLab → Build → Dev/운영 Pipeline, eCAMS, 취약점점검", "CI/CD"),
    ]
    table = slide.shapes.add_table(len(rows), 3, Inches(0.4), Inches(0.65), Inches(12.5), Inches(4.6)).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(7.2)
    table.columns[2].width = Inches(3.3)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)

    # I/F detail table
    add_label(slide, 0.35, 5.45, 12, 0.28, "I/F Layer 상세", 14, NAVY, True)
    if_rows = [
        ("구분", "항목"),
        ("공통 I/F", "SSO / Cruz APIM / EAI / FOS / MFT / ETL / CDC / 대내·대외 MCA / GSE"),
        ("관제·메타", "통합관제 / IT메타 / Nexus / Wily Collector / 배치작업관리"),
        ("연계·대외", "연계포털 / 빅데이터 / 계정계 / 대외기관 / 마이데이터 / 올원뱅크 / 디지털플랫폼 / 타금융"),
    ]
    t2 = slide.shapes.add_table(len(if_rows), 2, Inches(0.4), Inches(5.80), Inches(12.5), Inches(1.35)).table
    t2.columns[0].width = Inches(2.0)
    t2.columns[1].width = Inches(10.5)
    for r, row in enumerate(if_rows):
        for c, val in enumerate(row):
            t2.cell(r, c).text = val
    style_table(t2)


def build_slide4(prs):
    """형상배포 상세 — 도형 플로우."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_label(slide, 0.35, 0.15, 12, 0.35, "형상배포 Layer 상세", 20, NAVY, True)

    add_label(slide, 0.4, 0.65, 6, 0.28, "개발 경로", 14, NAVY, True)
    steps = [
        (0.4, "GitLab\n코드버전관리"),
        (2.3, "GitLabRunner\n빌드"),
        (4.2, "Dev Pipeline\n(+취약점점검)"),
        (6.3, "개발환경"),
    ]
    for i, (x, txt) in enumerate(steps):
        add_box(slide, x, 1.05, 1.7, 0.70, txt, fill=BLUE_LT, font=10, color=BLACK, bold=True)
        if i < len(steps) - 1:
            add_arrow_h(slide, x + 1.75, 1.30, 0.35, 0.18)
            if i == 0:
                add_label(slide, x + 1.70, 1.55, 0.8, 0.2, "dev merge", 9, RED, True, PP_ALIGN.CENTER)

    add_label(slide, 0.4, 2.20, 10, 0.28, "운영 경로 (Master merge · 통제점검)", 14, NAVY, True)
    add_box(slide, 0.4, 2.65, 1.5, 0.55, "eCAMS", fill=GRAY_BOX, font=11, color=BLACK, bold=True)
    add_arrow_h(slide, 2.0, 2.82, 0.35, 0.18)
    add_box(slide, 2.45, 2.65, 1.5, 0.55, "빌드\n(+취약점점검)", fill=BLUE_LT, font=10, color=BLACK, bold=True)
    add_arrow_h(slide, 4.05, 2.82, 0.35, 0.18)
    add_box(slide, 4.50, 2.65, 1.5, 0.55, "운영환경", fill=GREEN, font=11, color=WHITE, bold=True)

    add_box(slide, 0.4, 3.50, 1.5, 0.55, "빌드\nGitLabRunner", fill=BLUE_LT, font=10, color=BLACK, bold=True)
    add_arrow_h(slide, 2.0, 3.67, 0.35, 0.18)
    add_box(slide, 2.45, 3.50, 1.7, 0.55, "운영-Pipeline", fill=BLUE_LT, font=11, color=BLACK, bold=True)
    add_arrow_h(slide, 4.25, 3.67, 0.35, 0.18)
    add_box(slide, 4.70, 3.50, 1.5, 0.55, "운영환경", fill=GREEN, font=11, color=WHITE, bold=True)

    # legend table
    add_label(slide, 0.4, 4.40, 8, 0.28, "범례", 14, NAVY, True)
    legend = [
        ("도형/색", "의미"),
        ("파란 WAS/WEB", "서비스 런타임 노드"),
        ("초록 IMDG / 운영환경", "캐시·운영 반영 지점"),
        ("빨간 텍스트", "싱글뷰 / merge / Session Cache 등 핵심 호출"),
        ("회색 I/F 박스", "공통·대외 인터페이스"),
    ]
    t = slide.shapes.add_table(len(legend), 2, Inches(0.4), Inches(4.80), Inches(8.5), Inches(2.0)).table
    t.columns[0].width = Inches(2.8)
    t.columns[1].width = Inches(5.7)
    for r, row in enumerate(legend):
        for c, val in enumerate(row):
            t.cell(r, c).text = val
    style_table(t)


def main():
    content_path = prepare_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide1(prs)
    build_slide2(prs, content_path)
    build_slide3(prs)
    build_slide4(prs)

    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    main()
